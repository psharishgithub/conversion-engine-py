from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from typing import List
import os
import re
import asyncpg
from openai import OpenAI
import numpy as np
from fastapi.middleware.cors import CORSMiddleware
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from fastapi.responses import StreamingResponse
import json
import asyncio

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],  # Adjust this to match your frontend URL
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Set the environment (development or production)
ENVIRONMENT = 'development'

# Set temporary directory based on environment
# TEMP_DIR = '/tmp' if ENVIRONMENT == 'production' else 'C:\\Users\\Harish\\AppData\\Local\\Temp'
TEMP_DIR = 'C:\\Users\\Harish\\AppData\\Local\\Temp'

# Initialize OpenAI API client
client = OpenAI(
    base_url="https://laptop-483nic2i.tail7526d.ts.net/v1",
    api_key="lm-studio"
)

# Database connection settings
DATABASE_URL = "postgresql://postgres:6969@localhost:5432/recrepodb"

# Define the request model for file processing
class ProcessRequest(BaseModel):
    subject_id: str
    subject_code: str
    filepaths: List[str]
    user_id: str

# Define the request model for chat
class ChatRequest(BaseModel):
    question: str
    subject_id: str

@app.get("/")
async def root():
    return {"message": "hi"}

def extract_from_file(filepath: str) -> str:
    full_path = os.path.join(TEMP_DIR, os.path.basename(filepath))
    try:
        if (full_path.lower().endswith('.pdf')):
            print(f"Processing supported file: {full_path}")
            return extract_text_from_pdf(full_path)
        elif (full_path.lower().endswith('.docx')):
            print(f"Processing supported file: {full_path}")
            return extract_text_from_docx(full_path)
        elif (full_path.lower().endswith('.pptx')):
            print(f"Processing supported file: {full_path}")
            return extract_text_from_pptx(full_path)
        else:
            print(f"Unsupported file type: {full_path}")
            return f"Unsupported file type: {full_path}\n"
    except Exception as e:
        print(f"Error processing file {full_path}: {e}")
        return f"Error processing file {full_path}: {e}\n"

def extract_text_from_pdf(filepath: str) -> str:
    try:
        with open(filepath, 'rb') as file:
            reader = PdfReader(file)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            return '\n'.join(text)
    except Exception as e:
        print(f"Error processing PDF file {filepath}: {e}")
        return f"Error processing PDF file {filepath}: {e}\n"

def extract_text_from_docx(filepath: str) -> str:
    try:
        doc = Document(filepath)
        text = [para.text for para in doc.paragraphs]
        return '\n'.join(text)
    except Exception as e:
        print(f"Error processing DOCX file {filepath}: {e}")
        return f"Error processing DOCX file {filepath}: {e}\n"

def extract_text_from_pptx(filepath: str) -> str:
    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error processing PPTX file {filepath}: {e}")
        return f"Error processing PPTX file {filepath}: {e}\n"

# Process and clean text
def process_text(content: str) -> str:
    content = re.sub(r'\n+', '\n', content)
    content = re.sub(r'\s+', ' ', content)
    content = '\n'.join(line.strip() for line in content.split('\n'))
    return content

# Split text into chunks
def split_into_chunks(text: str, max_chunk_size: int = 2000) -> List[str]:
    chunks = []
    current_chunk = []
    words = text.split()
    for word in words:
        if len(' '.join(current_chunk) + ' ' + word) > max_chunk_size:
            chunks.append(' '.join(current_chunk))
            current_chunk = [word]
        else:
            current_chunk.append(word)
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    return chunks

# Get embeddings from OpenAI API
async def get_embedding(text: str, model: str = "nomic-ai/nomic-embed-text-v1.5-GGUF") -> List[float]:
    text = text.replace("\n", " ")
    response = client.embeddings.create(input=[text], model=model)
    return response.data[0].embedding

# Connect to PostgreSQL
async def connect_db():
    return await asyncpg.connect(DATABASE_URL)

# Upload aggregated text and embeddings to PostgreSQL
async def upload_to_db(subject_id: str, subject_code: str, aggregated_content: str, chunks: List[str], embeddings: List[dict], conn):
    try:
        async with conn.transaction():
            # Check if the subject exists
            subject = await conn.fetchrow('SELECT * FROM "Subject" WHERE "id" = $1', subject_id)
            
            if subject:
                # Update existing subject
                await conn.execute('''
                    UPDATE "Subject"
                    SET "aggregatedContent" = $2
                    WHERE "id" = $1;
                ''', subject_id, aggregated_content)
            else:
                # Insert new subject
                await conn.execute('''
                    INSERT INTO "Subject" ("id", "code", "aggregatedContent")
                    VALUES ($1, $2, $3);
                ''', subject_id, subject_code, aggregated_content)
            
            # Insert or update embeddings
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                await conn.execute('''
                    INSERT INTO "Embedding" ("subjectId", "chunkId", "text", "vector")
                    VALUES ($1, $2, $3, $4)
                    ON CONFLICT ("subjectId", "chunkId") 
                    DO UPDATE SET 
                        "text" = EXCLUDED."text", 
                        "vector" = EXCLUDED."vector";
                ''', subject_id, i, chunk, embedding['embedding'])
    except asyncpg.exceptions.PostgresError as e:
        print(f"Database error: {e}")
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(f"Unexpected error: {e}")
        raise HTTPException(status_code=500, detail="An unexpected error occurred")

@app.post("/process")
async def process(request: ProcessRequest):
    if not request.subject_code:
        raise HTTPException(status_code=400, detail="Subject code is required")

    content = ''
    for filepath in request.filepaths:
        full_path = os.path.join(TEMP_DIR, os.path.basename(filepath))
        if os.path.isfile(full_path):
            content += extract_from_file(filepath) + "\n"  
        else:
            content += f"File not found: {full_path}\n"
    
    textProcessed = process_text(content)
    
    if not textProcessed.strip():
        raise HTTPException(status_code=400, detail="No valid content extracted from files")
    
    chunks = split_into_chunks(textProcessed)
    
    # Convert chunks to embeddings
    embeddings = []
    for chunk in chunks:
        try:
            embedding = await get_embedding(chunk)
            embeddings.append({
                "text": chunk,
                "embedding": embedding
            })
        except Exception as e:
            print(f"Error getting embedding for chunk: {chunk[:30]}...: {e}")
            raise HTTPException(status_code=500, detail="Error getting embeddings")

    # Connect to the database
    conn = await connect_db()
    
    try:
        # Call the upload function to update the database
        await upload_to_db(request.subject_id, request.subject_code, textProcessed, chunks, embeddings, conn)
    except Exception as e:
        print(f"Error uploading to database: {e}")
        raise HTTPException(status_code=500, detail="Error uploading to database")
    finally:
        await conn.close()

    return {
        "message": "Files processed successfully",
        "aggregated_content": textProcessed,
        "chunks": chunks,
        "embeddings": embeddings
    }

@app.post("/chat")
async def chat(request: ChatRequest):
    if not request.question:
        raise HTTPException(status_code=400, detail="Question is required")

    async def generate_response():
        conn = await connect_db()
        try:
            # Check if it's a greeting
            greetings = ['hi', 'hello', 'hey', 'hi there', 'hello there']
            is_greeting = request.question.lower().strip() in greetings

            if is_greeting:
                greeting_response = {
                    "content": "Hello! I'm your AI teaching assistant. I'm here to help you understand your subject materials better. What topic would you like to explore today?"
                }
                yield json.dumps(greeting_response) + "\n"
                return

            # Rest of the embedding retrieval logic
            embeddings = await conn.fetch('''
                SELECT "chunkId", "text", "vector"
                FROM "Embedding"
                WHERE "subjectId" = $1
            ''', request.subject_id)
            
            if not embeddings:
                yield json.dumps({"error": "No embeddings found"}) + "\n"
                return

            question_embedding = await get_embedding(request.question)
            
            # Find best chunk (same logic as before)
            best_chunk = None
            best_similarity = -1

            for embedding in embeddings:
                chunk_embedding = np.array(embedding['vector'])
                similarity = np.dot(question_embedding, chunk_embedding) / (np.linalg.norm(question_embedding) * np.linalg.norm(chunk_embedding))
                
                if similarity > best_similarity:
                    best_similarity = similarity
                    best_chunk = embedding['text']

            if best_chunk is None:
                yield json.dumps({"error": "No relevant content found"}) + "\n"
                return

            # Enhanced context-aware prompt with better formatting instructions
            prompt = f"""As a teaching assistant, help the student understand this topic. Structure your response following these specific formatting rules:

            1. Start with a clear main heading using # (Example: # Topic Title)
            2. Use ## for major subtopics with double line breaks before and after
            3. Use ### for smaller subtopics with single line breaks
            4. Use bold (**text**) for important terms and concepts
            5. Use bullet points with proper indentation
            6. Use `code blocks` for specific terms, formulas, or technical details
            7. Add horizontal lines (---) between major sections

            Use this reference material to inform your response:
            {best_chunk}

            Student's question: {request.question}

            Format your response using this structure:
            
            # Main Topic
            
            Brief introduction
            
            ## Key Concepts
            
            * **Important term**: Definition or explanation
            * **Another term**: Description
            
            ## Detailed Explanation
            
            ### Subtopic 1
            Explanation with `specific terms` highlighted
            
            ### Subtopic 2
            More details with examples
            
            ---
            
            ## Summary
            Concise recap of main points

            Response:"""

            # Stream the completion with enhanced role
            stream = client.chat.completions.create(
                model="lmstudio-community/Meta-Llama-3-8B-Instruct-GGUF",
                messages=[
                    {
                        "role": "system", 
                        "content": """You are an experienced teaching assistant who creates well-structured, visually organized notes. 
                        Always use proper Markdown formatting with consistent spacing and highlighting of key concepts. 
                        Make sure to separate sections with appropriate line breaks and use emphasis markers for important terms."""
                    },
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                stream=True
            )

            for chunk in stream:
                if chunk.choices[0].delta.content is not None:
                    yield json.dumps({"content": chunk.choices[0].delta.content}) + "\n"
                    await asyncio.sleep(0.05)

        except Exception as e:
            yield json.dumps({"error": str(e)}) + "\n"
        finally:
            await conn.close()

    return StreamingResponse(
        generate_response(),
        media_type='text/event-stream',
        headers={
            'Content-Type': 'text/event-stream',
            'Cache-Control': 'no-cache',
            'Connection': 'keep-alive',
        }
    )

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)